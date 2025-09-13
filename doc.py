import streamlit as st
import PyPDF2
from io import BytesIO
from gtts import gTTS
import tempfile
import os
import docx  # for Word
from pptx import Presentation  # for PowerPoint
import boto3

# ---------------------------
# Custom CSS Styling
# ---------------------------
st.markdown("""
<style>
    /* Main app background and font */
    .main {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    
    /* Custom container styling */
    .stApp {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    }
    
    /* Title styling */
    .main-title {
        text-align: center;
        color: #ffffff;
        font-size: 3rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
    }
    
    .subtitle {
        text-align: center;
        color: #e8f4f8;
        font-size: 1.2rem;
        margin-bottom: 2rem;
        text-shadow: 1px 1px 2px rgba(0,0,0,0.2);
    }
    
    /* Card-like containers */
    .custom-card {
        background: rgba(255, 255, 255, 0.95);
        backdrop-filter: blur(10px);
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: 0 8px 32px rgba(0,0,0,0.1);
        border: 1px solid rgba(255, 255, 255, 0.2);
    }
    
    /* Section headers */
    .section-header {
        color: #2c3e50;
        font-size: 1.5rem;
        font-weight: 600;
        margin-bottom: 1rem;
        border-bottom: 3px solid #667eea;
        padding-bottom: 0.5rem;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    /* File upload area styling */
    .uploadedFile {
        background: #f8f9fa;
        border: 2px dashed #667eea;
        border-radius: 10px;
        padding: 1rem;
        text-align: center;
        margin: 1rem 0;
    }
    
    /* Button styling */
    .stButton > button {
        background: linear-gradient(45deg, #667eea, #764ba2);
        color: white;
        border: none;
        border-radius: 25px;
        padding: 0.75rem 2rem;
        font-size: 1rem;
        font-weight: 600;
        cursor: pointer;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.3);
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.4);
    }
    
    /* Selectbox styling */
    .stSelectbox > div > div {
        background: white;
        color:black;
        border-radius: 10px;
        border: 2px solid #e9ecef;
    }
    
    .stSelectbox > div > div:focus-within {
        border-color: #667eea;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.1);
    }
    
    /* Text area styling */
    .stTextArea > div > div > textarea {
        background: #f8f9fa;
        border: 2px solid #e9ecef;
        border-radius: 10px;
        font-family: 'Consolas', 'Monaco', monospace;
        font-size: 0.9rem;
        line-height: 1.5;
    }
    
    /* Success/Error message styling */
    .stSuccess {
        background: linear-gradient(45deg, #28a745, #20c997);
        border-radius: 10px;
        color: white;
    }
    
    .stError {
        background: linear-gradient(45deg, #dc3545, #e83e8c);
        border-radius: 10px;
        color: white;
    }
    
    /* Audio player styling */
    audio {
        width: 100%;
        border-radius: 10px;
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background: rgba(255, 255, 255, 0.95);
        backdrop-filter: blur(10px);
    }
    
    /* Info boxes */
    .info-box {
        background: linear-gradient(135deg, #74b9ff, #0984e3);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        margin: 1rem 0;
        text-align: center;
        box-shadow: 0 4px 15px rgba(116, 185, 255, 0.3);
    }
    
    /* Spinner customization */
    .stSpinner {
        text-align: center;
        color: #667eea;
    }
    
    /* File uploader styling */
    .stFileUploader > div {
        border: 2px dashed #667eea;
        border-radius: 15px;
        padding: 2rem;
        background: rgba(255, 255, 255, 0.8);
        transition: all 0.3s ease;
    }
    
    .stFileUploader > div:hover {
        border-color: #764ba2;
        background: rgba(255, 255, 255, 0.9);
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background: linear-gradient(45deg, #667eea, #764ba2);
    }
    
    /* Link styling */
    a {
        color: #667eea;
        text-decoration: none;
        font-weight: 600;
    }
    
    a:hover {
        color: #764ba2;
        text-decoration: underline;
    }
    
    /* Responsive design */
    @media (max-width: 768px) {
        .main-title {
            font-size: 2rem;
        }
        .custom-card {
            padding: 1rem;
            margin: 0.5rem 0;
        }
    }
</style>
""", unsafe_allow_html=True)

# ---------------------------
# Initialize S3 Client
# ---------------------------
s3 = boto3.client('s3')

def generate_audio(text, lang="en"):
    """Generate MP3 from text using gTTS."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp_file:
        tts = gTTS(text=text, lang=lang, slow=False)
        tts.save(tmp_file.name)
        with open(tmp_file.name, "rb") as audio_file:
            audio_bytes = audio_file.read()
        os.unlink(tmp_file.name)
    return audio_bytes, "output.mp3"

def upload_to_s3(file_path, bucket_name, object_name):
    """Upload MP3 to S3 and return shareable link."""
    s3.upload_file(file_path, bucket_name, object_name, ExtraArgs={'ACL': 'public-read'})
    return f"https://{bucket_name}.s3.amazonaws.com/{object_name}"

# ---------------------------
# Streamlit UI
# ---------------------------
st.set_page_config(
    page_title="File to Speech Converter", 
    page_icon="📖", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom title with styling
st.markdown('<h1 class="main-title">📖 File to Speech Converter</h1>', unsafe_allow_html=True)
st.markdown('<p class="subtitle">Transform your documents into audio with AI-powered speech synthesis</p>', unsafe_allow_html=True)

# --- Session State Initialization ---
if "extracted_text" not in st.session_state:
    st.session_state["extracted_text"] = ""
if "audio_bytes" not in st.session_state:
    st.session_state["audio_bytes"] = None
if "share_link" not in st.session_state:
    st.session_state["share_link"] = None

# ---------------------------
# Extraction Functions (cached)
# ---------------------------
@st.cache_data
def extract_text_from_pdf(file_bytes):
    reader = PyPDF2.PdfReader(BytesIO(file_bytes))
    texts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(texts)

@st.cache_data
def extract_text_from_docx(file_bytes):
    doc = docx.Document(BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

@st.cache_data
def extract_text_from_pptx(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

@st.cache_data
def extract_text_from_txt(file_bytes):
    return BytesIO(file_bytes).read().decode("utf-8", errors="ignore")

# ---------------------------
# Main Content Area
# ---------------------------
col1, col2 = st.columns([2, 1])

with col1:
    # File Upload Section
    st.markdown('<div class="custom-card">', unsafe_allow_html=True)
    st.markdown('<h2 class="section-header">📁 File Upload</h2>', unsafe_allow_html=True)
    
    file_type = st.selectbox(
        "Select File Type", 
        ["PDF", "Word Document", "PowerPoint", "Text File"],
        help="Choose the type of document you want to convert to speech"
    )

    if file_type == "PDF":
        uploaded_file = st.file_uploader("Choose a PDF file", type=["pdf"])
    elif file_type == "Word Document":
        uploaded_file = st.file_uploader("Choose a Word file", type=["docx"])
    elif file_type == "PowerPoint":
        uploaded_file = st.file_uploader("Choose a PowerPoint file", type=["pptx"])
    elif file_type == "Text File":
        uploaded_file = st.file_uploader("Choose a Text file", type=["txt"])
    else:
        uploaded_file = None
    
    st.markdown('</div>', unsafe_allow_html=True)

    # File Processing Section
    if uploaded_file:
        try:
            file_bytes = uploaded_file.read()
            text = ""

            with st.spinner("📖 Processing your document..."):
                if uploaded_file.name.endswith(".pdf"):
                    text = extract_text_from_pdf(file_bytes)
                elif uploaded_file.name.endswith(".docx"):
                    text = extract_text_from_docx(file_bytes)
                elif uploaded_file.name.endswith(".pptx"):
                    text = extract_text_from_pptx(file_bytes)
                elif uploaded_file.name.endswith(".txt"):
                    text = extract_text_from_txt(file_bytes)

            if text.strip():
                st.session_state["extracted_text"] = text

                # Extracted Text Section
                st.markdown('<div class="custom-card">', unsafe_allow_html=True)
                st.markdown('<h2 class="section-header">📄 Extracted Text</h2>', unsafe_allow_html=True)
                
                # Show file info
                st.info(f"📊 **File:** {uploaded_file.name} | **Size:** {len(file_bytes):,} bytes | **Text Length:** {len(text):,} characters")
                
                st.text_area(
                    "Document Content:", 
                    text, 
                    height=300, 
                    disabled=True,
                    help="This is the extracted text from your document"
                )
                st.markdown('</div>', unsafe_allow_html=True)

                # Audio Generation Section
                st.markdown('<div class="custom-card">', unsafe_allow_html=True)
                st.markdown('<h2 class="section-header">🎙️ Audio Generation</h2>', unsafe_allow_html=True)
                
                col_lang, col_btn = st.columns([1, 1])
                
                with col_lang:
                    lang_options = {
                        "English": "en",
                        "Hindi": "hi", 
                        "Spanish": "es",
                        "French": "fr",
                        "German": "de",
                        "Italian": "it",
                        "Portuguese": "pt",
                        "Russian": "ru",
                        "Japanese": "ja",
                        "Korean": "ko"
                    }
                    selected_lang = st.selectbox(
                        "🌐 Select Language", 
                        options=list(lang_options.keys()),
                        help="Choose the language for text-to-speech conversion"
                    )
                    lang = lang_options[selected_lang]

                with col_btn:
                    st.markdown("<br>", unsafe_allow_html=True)  # Add spacing
                    if st.button("🎵 Generate Audio", use_container_width=True):
                        with st.spinner("🎧 Creating your audio file..."):
                            st.session_state["audio_bytes"], file_name = generate_audio(text, lang)

                            # Save temp file for S3
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp_file:
                                tts = gTTS(text=text, lang=lang, slow=False)
                                tts.save(tmp_file.name)
                                st.session_state["share_link"] = upload_to_s3(tmp_file.name, "my-audio-bucket", file_name)
                                os.unlink(tmp_file.name)
                
                st.markdown('</div>', unsafe_allow_html=True)

            else:
                st.error("❌ No readable text could be extracted from this file. Please try a different document.")

        except Exception as e:
            st.error(f"❌ Error processing file: {str(e)}")

with col2:
    # Audio Player Section
    if st.session_state["audio_bytes"]:
        st.markdown('<div class="custom-card">', unsafe_allow_html=True)
        st.markdown('<h2 class="section-header">🎧 Audio Player</h2>', unsafe_allow_html=True)
        
        st.audio(st.session_state["audio_bytes"], format="audio/mp3")
        st.success("🎉 Audio generated successfully!")
        
        if st.session_state["share_link"]:
            st.markdown(
                f'<div class="info-box">🔗 <a href="{st.session_state["share_link"]}" target="_blank">Download/Share Audio File</a></div>', 
                unsafe_allow_html=True
            )
        
        st.markdown('</div>', unsafe_allow_html=True)

# ---------------------------
# Enhanced Sidebar
# ---------------------------
with st.sidebar:
    st.markdown('<div class="custom-card">', unsafe_allow_html=True)
    st.markdown("## 📖 How to Use")
    st.markdown("""
    **Step-by-step guide:**
    
    1. 📁 **Upload** a supported file
    2. 👀 **Review** the extracted text
    3. 🌐 **Choose** your preferred language
    4. 🎵 **Generate** the audio file
    5. 🎧 **Play** or share the audio
    """)
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="custom-card">', unsafe_allow_html=True)
    st.markdown("## 📋 Supported Formats")
    st.markdown("""
    - 📄 **PDF** (.pdf) - Portable documents
    - 📝 **Word** (.docx) - Microsoft Word docs
    - 📊 **PowerPoint** (.pptx) - Presentations
    - 📃 **Text** (.txt) - Plain text files
    """)
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="custom-card">', unsafe_allow_html=True)
    st.markdown("## ⚡ Features")
    st.markdown("""
    - 🎛️ **Speed Control** - Adjust playback speed
    - 💾 **Download Options** - Save audio files
    - 🌍 **Multi-language** - 10+ language support
    - ☁️ **Cloud Storage** - Shareable links
    - 🔄 **Real-time Processing** - Fast conversion
    """)
    st.markdown('</div>', unsafe_allow_html=True)

    # Add some metrics if there's extracted text
    if st.session_state["extracted_text"]:
        st.markdown('<div class="custom-card">', unsafe_allow_html=True)
        st.markdown("## 📊 Document Stats")
        text = st.session_state["extracted_text"]
        
        col_a, col_b = st.columns(2)
        with col_a:
            st.metric("Characters", f"{len(text):,}")
            st.metric("Words", f"{len(text.split()):,}")
        with col_b:
            st.metric("Lines", f"{len(text.splitlines()):,}")
            # Estimate reading time (average 200 words per minute)
            reading_time = max(1, len(text.split()) // 200)
            st.metric("Est. Reading Time", f"{reading_time} min")
        
        st.markdown('</div>', unsafe_allow_html=True)

# Footer
st.markdown("---")
st.markdown(
    '<p style="text-align: center; color: rgba(255,255,255,0.7); font-size: 0.9rem;">Made with ❤️ using Streamlit • Text-to-Speech powered by Google TTS</p>', 
    unsafe_allow_html=True
)